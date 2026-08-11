#!/usr/bin/env python3
"""
qa_deck.py - mechanical pre-render checks on a guide deck.

Implements the assertions specified in
  TR - Documents/4. Guides/2. Quality assurance/12_qr_integrity.md

Every check here is deterministic, so an LLM is the wrong executor for them: the
answer cannot change between runs and there is nothing to weigh. The agent file
stays as the specification; this is the implementation, and it is cheap enough to
run before every render rather than only at QA time.

Two failure families, both invisible to a human proof-reader because the page
renders perfectly in every one of them:

  QR      A code carries TWO independent targets - the slug encoded in the image,
          and the a:hlinkClick a PDF reader follows on tap. They are stored
          separately and drift apart when an image is swapped or a picture is
          cloned. Three builds in a row shipped a defect of this shape.

  HEADER  Decks are cloned from siblings, so the running head is inherited. A
          slide left saying "ICELAND | 7-DAY RING ROAD" inside the 5-day guide is
          template residue that no reader is meant to see, and it survives every
          prose check because nobody reads the running head.

Usage:
  python scripts/qa_deck.py "<deck.pptx>" --qr-dir content/countries/iceland/qr
  python scripts/qa_deck.py "<deck.pptx>" --qr-dir <dir> --sweep      # + production
  python scripts/qa_deck.py "<deck.pptx>" --qr-dir <dir> --json out.json
"""
import argparse
import collections
import io
import json
import pathlib
import re
import subprocess
import sys

from pptx import Presentation
from pptx.util import Emu

BASE = "testedroutes.com/go/"
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
CONTROL_SLUG = "qa-control-slug-does-not-exist"


def walk(shapes):
    """slide.shapes does not recurse into groups, and callout boxes live in them."""
    for sh in shapes:
        if sh.shape_type == 6:
            yield from walk(sh.shapes)
        else:
            yield sh


# --------------------------------------------------------------------- QR checks
def qr_inventory(prs):
    out = []
    for i, slide in enumerate(prs.slides, 1):
        for sh in walk(slide.shapes):
            if sh.shape_type != 13:
                continue
            name = sh.name or ""
            if not name.startswith("QR"):
                continue
            alias = name[3:].strip() if name.startswith("QR ") else None
            try:
                link = sh.click_action.hyperlink.address
            except Exception:
                link = None
            out.append({"slide": i, "shape": name, "alias": alias, "hyperlink": link,
                        "blob": sh.image.blob if alias else None,
                        "top": round(Emu(sh.top).pt), "left": round(Emu(sh.left).pt)})
    return out


def payload_matches(blob, alias, qr_dir):
    """Compare the embedded image against the alias's PNG in the country QR pool.

    An earlier version re-encoded the expected payload and compared pixels. That
    is wrong: QR generators differ in version, module size and error-correction,
    so every code not produced by this exact script failed - 24 false criticals on
    a 31-code deck. The pool file IS the source of truth for what should be
    embedded, so compare against it and the check is parameter-agnostic.

    Returns True, False, or None when the pool has no file for that alias (which
    is itself reported, as orphan_alias)."""
    if not qr_dir:
        return None
    f = pathlib.Path(qr_dir) / (alias + ".png")
    if not f.exists():
        return None
    if blob == f.read_bytes():
        return True
    # Same code re-saved at a different size still scans identically; fall back to
    # a pixel comparison at a common size before calling it a mismatch.
    try:
        from PIL import Image, ImageChops
    except ImportError:
        return False
    a_img = Image.open(io.BytesIO(blob)).convert("L").resize((330, 330), Image.NEAREST)
    b_img = Image.open(f).convert("L").resize((330, 330), Image.NEAREST)
    return ImageChops.difference(a_img, b_img).getbbox() is None


def sweep(slugs):
    """Resolve every slug against production, with a control that must come back
    dead. A sweep whose control passes cannot detect failure, and its clean result
    means nothing - this has produced a false pass twice."""
    results = {}
    for slug in list(slugs) + [CONTROL_SLUG]:
        slug = slug.replace("\r", "").strip()      # a \r makes every URL wrong
        if not slug:
            continue
        try:
            url = subprocess.run(
                ["curl", "-s", "-o", "/dev/null", "-w", "%{redirect_url}",
                 "--max-time", "20", "https://testedroutes.com/go/" + slug],
                capture_output=True, text=True, timeout=30).stdout.strip()
        except Exception:
            url = ""
        results[slug] = url
    return results


# ----------------------------------------------------------------- header checks
HEAD_RE = re.compile(r"^[A-Z][A-Z\s]*\|\s*\S")


def header_inventory(prs, top_max=40):
    """Running heads sit in the top band. Collect every candidate so a stale one
    from a sibling deck is visible next to the majority."""
    heads = []
    for i, slide in enumerate(prs.slides, 1):
        for sh in walk(slide.shapes):
            if not sh.has_text_frame:
                continue
            t = " ".join(sh.text_frame.text.split())
            if not t or Emu(sh.top).pt > top_max:
                continue
            if HEAD_RE.match(t) and len(t) < 60:
                heads.append({"slide": i, "shape": sh.name, "text": t})
    return heads


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("--qr-dir", help="country QR pool, for orphan-PNG reporting")
    ap.add_argument("--expect-head", help="running head; inferred from the majority if omitted")
    ap.add_argument("--sweep", action="store_true", help="resolve slugs against production")
    ap.add_argument("--json", help="write the full report here")
    a = ap.parse_args()

    prs = Presentation(a.deck)
    findings, trace = [], []

    # ---- headers
    heads = header_inventory(prs)
    counts = collections.Counter(h["text"] for h in heads)
    expected = a.expect_head or (counts.most_common(1)[0][0] if counts else None)
    for h in heads:
        ok = h["text"] == expected
        trace.append({"check": "header", "slide": h["slide"], "value": h["text"],
                      "status": "pass" if ok else "fail"})
        if not ok:
            findings.append({"type": "header_mismatch", "severity": "major",
                             "slide": h["slide"], "location": h["text"],
                             "expected": expected,
                             "finding": "running head does not match the deck majority - "
                                        "usually a slide pasted in from a sibling guide",
                             "confidence": 1.0})
    missing_head = [i for i in range(1, len(prs.slides._sldIdLst) + 1)
                    if not any(h["slide"] == i for h in heads)]
    for i in missing_head:
        findings.append({"type": "header_missing", "severity": "minor", "slide": i,
                         "location": "(none)", "expected": expected,
                         "finding": "no running head found on this slide",
                         "confidence": 1.0})

    # ---- QRs
    inv = qr_inventory(prs)
    seen_pos = collections.defaultdict(list)
    for q in inv:
        if not q["alias"]:
            findings.append({"type": "unnamed_qr", "severity": "major",
                             "slide": q["slide"], "location": q["shape"],
                             "finding": "QR shape is not named after its alias, so nothing "
                                        "can ever check it", "confidence": 1.0})
            continue
        seen_pos[(q["slide"], q["top"], q["left"])].append(q["alias"])
        pm = payload_matches(q["blob"], q["alias"], a.qr_dir)
        hl = q["hyperlink"]
        hm = bool(hl) and ("go/" + q["alias"]) in hl
        trace.append({"check": "qr", "slide": q["slide"], "alias": q["alias"],
                      "payload_matches": pm, "hyperlink": hl,
                      "hyperlink_matches": hm,
                      "status": "pass" if (pm is not False and hm) else "fail"})
        if pm is None:
            findings.append({"type": "orphan_alias", "severity": "major",
                             "slide": q["slide"], "location": q["shape"],
                             "alias": q["alias"],
                             "finding": "no PNG for this alias in the country QR pool, so "
                                        "the embedded image cannot be verified",
                             "confidence": 1.0})
        if pm is False:
            findings.append({"type": "payload_mismatch", "severity": "critical",
                             "slide": q["slide"], "location": q["shape"],
                             "alias": q["alias"],
                             "finding": "encoded image does not match the shape's alias",
                             "confidence": 1.0})
        if not hl:
            findings.append({"type": "hyperlink_missing", "severity": "critical",
                             "slide": q["slide"], "location": q["shape"],
                             "alias": q["alias"],
                             "finding": "QR has no hyperlink - scanning works, tapping does "
                                        "nothing", "confidence": 1.0})
        elif not hm:
            findings.append({"type": "hyperlink_mismatch", "severity": "critical",
                             "slide": q["slide"], "location": q["shape"],
                             "alias": q["alias"], "hyperlink": hl,
                             "finding": "image encodes %s but the hyperlink opens %s"
                                        % (q["alias"], hl.rsplit("/", 1)[-1]),
                             "confidence": 1.0})
    for pos, aliases in seen_pos.items():
        if len(set(aliases)) > 1:
            findings.append({"type": "duplicate_position", "severity": "minor",
                             "slide": pos[0], "location": str(pos),
                             "finding": "two different aliases share one position: %s"
                                        % ", ".join(sorted(set(aliases))),
                             "confidence": 1.0})

    control = None
    if a.sweep:
        res = sweep({q["alias"] for q in inv if q["alias"]})
        control = {"slug": CONTROL_SLUG,
                   "detected_dead": "utm_campaign=unknown" in res.get(CONTROL_SLUG, "")}
        for slug, url in res.items():
            if slug == CONTROL_SLUG:
                continue
            dead = (not url) or "utm_campaign=unknown" in url
            trace.append({"check": "sweep", "alias": slug, "resolves_to": url,
                          "status": "fail" if dead else "pass"})
            if dead:
                findings.append({"type": "dead_slug", "severity": "critical",
                                 "location": slug, "alias": slug,
                                 "finding": "resolves to the homepage in production",
                                 "confidence": 1.0})
        if not control["detected_dead"]:
            findings.append({"type": "control_failed", "severity": "critical",
                             "location": CONTROL_SLUG,
                             "finding": "control slug did not come back dead - every "
                                        "dead_slug result in this run is unreliable",
                             "confidence": 1.0})

    report = {"deck": str(pathlib.Path(a.deck).name),
              "counts": {"slides": len(prs.slides._sldIdLst), "qr": len(inv),
                         "headers": len(heads), "findings": len(findings)},
              "expected_head": expected, "control": control,
              "trace": trace, "findings": findings}

    crit = sum(1 for f in findings if f["severity"] == "critical")
    print("deck        : %s" % report["deck"])
    print("slides      : %d   QRs: %d   headers: %d"
          % (report["counts"]["slides"], len(inv), len(heads)))
    print("running head: %s" % expected)
    for f in findings:
        print("  [%-8s] s%-2s %-22s %s"
              % (f["severity"], f.get("slide", "-"), f["type"], f["finding"][:78]))
    print("findings    : %d (%d critical)" % (len(findings), crit))
    if a.json:
        pathlib.Path(a.json).write_text(json.dumps(report, indent=1), encoding="utf-8")
        print("report      : %s" % a.json)
    return 1 if crit else 0


if __name__ == "__main__":
    sys.exit(main())
